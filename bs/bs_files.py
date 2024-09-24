import requests
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from fpdf import FPDF

def extract_research_info(latex_code):
    """Extract research information based on the given LaTeX code."""
    search_query = latex_code.replace("$", "").replace("\\", "").replace("{", "").replace("}", "")
    url = f"https://www.google.com/search?q={search_query}"

    headers = {'User-Agent': 'Mozilla/5.0'}
    response = requests.get(url, headers=headers)
    soup = BeautifulSoup(response.text, 'html.parser')

    results = []
    for g in soup.find_all('div', class_='BVG0Nb'):
        title = g.find('h3').text
        snippet = g.find('div', class_='BNeawe').text
        results.append(f"{title}: {snippet}")
    
    return results  # Return a list of research snippets

def to_word(latex_code, research_info, filename='output.docx'):
    """Convert LaTeX code and research info to a Word document."""
    doc = Document()
    doc.add_paragraph('Mathematical Equation:')
    doc.add_paragraph(latex_code)  # Placeholder for rendering LaTeX

    for info in research_info:
        doc.add_page_break()  # Add a page break for each research info
        doc.add_paragraph('Research Information:')
        doc.add_paragraph(info)  # Add each research info entry

    doc.save(filename)

def to_ppt(latex_code, research_info, filename='output.pptx'):
    """Convert LaTeX code and research info to a PowerPoint presentation."""
    prs = Presentation()

    # Add a slide for the LaTeX code
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use a blank slide
    title = slide.shapes.title
    title.text = "LaTeX Equation"
    
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(3))
    text_frame = textbox.text_frame
    text_frame.text = latex_code  # Placeholder for rendering LaTeX

    # Add a new slide for each research info entry
    for info in research_info:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use a blank slide
        title = slide.shapes.title
        title.text = "Research Information"
        
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
        text_frame = textbox.text_frame
        text_frame.text = info  # Add research info

    prs.save(filename)

def to_pdf(latex_code, research_info, filename='output.pdf'):
    """Convert LaTeX code and research info to a PDF file."""
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)

    pdf.add_page()
    pdf.set_font("Arial", size=12)
    pdf.cell(200, 10, txt="Mathematical Equation:", ln=True, align='C')
    pdf.cell(200, 10, txt=latex_code, ln=True, align='C')

    for info in research_info:
        pdf.add_page()  # Add a new page for each research info
        pdf.cell(200, 10, txt="Research Information:", ln=True, align='C')
        pdf.multi_cell(0, 10, txt=info)  # Allow multi-line for research info

    pdf.output(filename)
