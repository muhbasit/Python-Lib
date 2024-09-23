# staT.py
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from openpyxl import Workbook
from fpdf import FPDF
import matplotlib.pyplot as plt 

class staT:
    """A class to perform basic statistical operations on both raw and grouped data, with working tables."""

    @staticmethod
    def mean(data, grouped=False):
        """Calculate the mean of the data."""
        if grouped:
            total_frequency = sum(freq for _, freq in data)
            total_sum = sum(midpoint * freq for midpoint, freq in data)
            return total_sum / total_frequency if total_frequency > 0 else 0
        else:
            total = sum(data)
            count = len(data)
            return total / count if count > 0 else 0

    @staticmethod
    def median(data, grouped=False):
        """Calculate the median of the data."""
        if grouped:
            cumulative_frequency = 0
            total_frequency = sum(freq for _, freq in data)
            for midpoint, freq in data:
                cumulative_frequency += freq
                if cumulative_frequency >= total_frequency / 2:
                    return midpoint
        else:
            sorted_data = sorted(data)
            count = len(sorted_data)
            if count % 2 == 0:
                mid1 = count // 2
                mid2 = mid1 - 1
                return (sorted_data[mid1] + sorted_data[mid2]) / 2
            else:
                mid = count // 2
                return sorted_data[mid]

    @staticmethod
    def mode(data, grouped=False):
        """Calculate the mode of the data."""
        if grouped:
            max_frequency = max(freq for _, freq in data)
            mode_values = [midpoint for midpoint, freq in data if freq == max_frequency]
            return mode_values if max_frequency > 0 else None
        else:
            frequency = {}
            for value in data:
                frequency[value] = frequency.get(value, 0) + 1
            
            max_count = max(frequency.values())
            mode_values = [key for key, count in frequency.items() if count == max_count]
            return mode_values if max_count > 1 else None  # return None if there is no mode

    @staticmethod
    def mct(data, grouped=False):
        """Calculate measures of central tendency: mean, median, and mode."""
        mean_value = staT.mean(data, grouped)
        median_value = staT.median(data, grouped)
        mode_value = staT.mode(data, grouped)
        return [mean_value, median_value, mode_value]

    @staticmethod
    def mct_wr(data, grouped=False, return_all=True, return_mean=False, return_median=False, return_mode=False):
        """
        Returns working tables for calculating mean, median, and mode.
        
        Parameters:
        - return_all: If True, returns working tables for mean, median, and mode.
        - return_mean: If True, returns the working table for mean.
        - return_median: If True, returns the working table for median.
        - return_mode: If True, returns the working table for mode.
        """
        working_tables = {}

        # Mean working table
        if return_all or return_mean:
            if grouped:
                total_frequency = sum(freq for _, freq in data)
                total_sum = sum(midpoint * freq for midpoint, freq in data)
                working_tables['mean'] = {
                    'midpoints': [midpoint for midpoint, _ in data],
                    'frequencies': [freq for _, freq in data],
                    'fx': [midpoint * freq for midpoint, freq in data],
                    'total_fx': total_sum,
                    'total_frequency': total_frequency,
                    'mean': total_sum / total_frequency if total_frequency > 0 else 0
                }
            else:
                working_tables['mean'] = {
                    'data': data,
                    'sum': sum(data),
                    'count': len(data),
                    'mean': sum(data) / len(data) if len(data) > 0 else 0
                }

        # Median working table
        if return_all or return_median:
            if grouped:
                cumulative_frequency = 0
                total_frequency = sum(freq for _, freq in data)
                working_tables['median'] = {
                    'midpoints': [midpoint for midpoint, _ in data],
                    'frequencies': [freq for _, freq in data],
                    'cumulative_frequencies': [],
                    'median': None
                }
                for midpoint, freq in data:
                    cumulative_frequency += freq
                    working_tables['median']['cumulative_frequencies'].append(cumulative_frequency)
                    if cumulative_frequency >= total_frequency / 2 and working_tables['median']['median'] is None:
                        working_tables['median']['median'] = midpoint
            else:
                sorted_data = sorted(data)
                count = len(sorted_data)
                if count % 2 == 0:
                    mid1 = count // 2
                    mid2 = mid1 - 1
                    median_value = (sorted_data[mid1] + sorted_data[mid2]) / 2
                else:
                    mid = count // 2
                    median_value = sorted_data[mid]
                working_tables['median'] = {
                    'sorted_data': sorted_data,
                    'median': median_value
                }

        # Mode working table
        if return_all or return_mode:
            if grouped:
                max_frequency = max(freq for _, freq in data)
                mode_values = [midpoint for midpoint, freq in data if freq == max_frequency]
                working_tables['mode'] = {
                    'midpoints': [midpoint for midpoint, _ in data],
                    'frequencies': [freq for _, freq in data],
                    'mode': mode_values if max_frequency > 0 else None
                }
            else:
                frequency = {}
                for value in data:
                    frequency[value] = frequency.get(value, 0) + 1
                max_count = max(frequency.values())
                mode_values = [key for key, count in frequency.items() if count == max_count]
                working_tables['mode'] = {
                    'data': data,
                    'frequencies': frequency,
                    'mode': mode_values if max_count > 1 else None
                }

        return working_tables

    @staticmethod
    def mct_formulae(grouped=False, return_all=True, return_mean=False, return_median=False, return_mode=False):
        """
        Returns the formulae for calculating mean, median, and mode along with detailed descriptions.
        
        Parameters:
        - grouped: Boolean, if True, returns formulae for grouped data. Default is False (ungrouped data).
        - return_all: Boolean, if True, returns formulae for mean, median, and mode. Default is True.
        - return_mean: Boolean, if True, returns the formula for mean only.
        - return_median: Boolean, if True, returns the formula for median only.
        - return_mode: Boolean, if True, returns the formula for mode only.
        
        Returns:
        - A dictionary containing the formulae and their descriptions for the requested measures.
        """
        formulae = {}

        # Mean formula with description
        if return_all or return_mean:
            if grouped:
                formulae['mean'] = {
                    'formula': "Mean (Grouped) = Σ(f_i * x_i) / Σ(f_i)",
                    'description': "The mean for grouped data is calculated as the sum of the product of each class midpoint (x_i) "
                                   "and its corresponding frequency (f_i), divided by the total sum of frequencies (Σ(f_i)).",
                    'parameters': {
                        'f_i': "Frequency of the i-th class",
                        'x_i': "Midpoint of the i-th class",
                        'Σ(f_i)': "Total sum of frequencies"
                    }
                }
            else:
                formulae['mean'] = {
                    'formula': "Mean (Ungrouped) = Σ(x_i) / n",
                    'description': "The mean for ungrouped data is calculated by summing all data points (x_i) and dividing by the number of data points (n).",
                    'parameters': {
                        'x_i': "The i-th data point",
                        'n': "Total number of data points"
                    }
                }

        # Median formula with description
        if return_all or return_median:
            if grouped:
                formulae['median'] = {
                    'formula': "Median (Grouped) = L + [(N/2 - CF) / f_m] * c",
                    'description': "The median for grouped data is calculated using the lower class boundary (L) of the median class, the cumulative frequency (CF) before the median class, "
                                   "the total frequency (N), the frequency of the median class (f_m), and the class width (c).",
                    'parameters': {
                        'L': "Lower boundary of the median class",
                        'N': "Total frequency (Σ(f_i))",
                        'CF': "Cumulative frequency before the median class",
                        'f_m': "Frequency of the median class",
                        'c': "Class width"
                    }
                }
            else:
                formulae['median'] = {
                    'formula': "Median (Ungrouped) = Middle value or (x_{n/2} + x_{(n/2 + 1)}) / 2",
                    'description': "For ungrouped data, if the dataset size is odd, the median is the middle value. If the size is even, the median is the average of the two middle values.",
                    'parameters': {
                        'n': "Total number of data points",
                        'x_{n/2}': "Middle value(s) when sorted"
                    }
                }

        # Mode formula with description
        if return_all or return_mode:
            if grouped:
                formulae['mode'] = {
                    'formula': "Mode (Grouped) = L + [(f_m - f_1) / (2f_m - f_1 - f_2)] * c",
                    'description': "The mode for grouped data is calculated using the frequency of the modal class (f_m), the frequencies of the classes before (f_1) and after (f_2) the modal class, "
                                   "the lower boundary of the modal class (L), and the class width (c).",
                    'parameters': {
                        'L': "Lower boundary of the modal class",
                        'f_m': "Frequency of the modal class",
                        'f_1': "Frequency of the class before the modal class",
                        'f_2': "Frequency of the class after the modal class",
                        'c': "Class width"
                    }
                }
            else:
                formulae['mode'] = {
                    'formula': "Mode (Ungrouped) = Most frequent value(s)",
                    'description': "The mode for ungrouped data is the value that appears most frequently in the dataset.",
                    'parameters': {
                        'Most frequent value': "Value(s) that occur most often in the dataset"
                    }
                }

        return formulae

    @staticmethod
    def mct_latex(grouped=False, return_all=True, return_mean=False, return_median=False, return_mode=False):
        """
        Returns LaTeX versions of the formulas for calculating mean, median, and mode.
        
        Parameters:
        - grouped: Boolean, if True, returns LaTeX formulas for grouped data. Default is False (ungrouped data).
        - return_all: Boolean, if True, returns LaTeX formulas for mean, median, and mode. Default is True.
        - return_mean: Boolean, if True, returns the LaTeX formula for mean only.
        - return_median: Boolean, if True, returns the LaTeX formula for median only.
        - return_mode: Boolean, if True, returns the LaTeX formula for mode only.
        
        Returns:
        - A dictionary containing the LaTeX versions of the requested formulas.
        """
        latex_formulas = {}

        # Mean LaTeX formula
        if return_all or return_mean:
            if grouped:
                latex_formulas['mean'] = r"\bar{x} = \frac{\sum (f_i \cdot x_i)}{\sum f_i}"
            else:
                latex_formulas['mean'] = r"\bar{x} = \frac{\sum x_i}{n}"

        # Median LaTeX formula
        if return_all or return_median:
            if grouped:
                latex_formulas['median'] = r"Median = L + \left( \frac{\frac{N}{2} - CF}{f_m} \right) \cdot c"
            else:
                latex_formulas['median'] = r"Median = \text{Middle value or} \frac{x_{\frac{n}{2}} + x_{\frac{n}{2}+1}}{2}"

        # Mode LaTeX formula
        if return_all or return_mode:
            if grouped:
                latex_formulas['mode'] = r"Mode = L + \left( \frac{f_m - f_1}{2f_m - f_1 - f_2} \right) \cdot c"
            else:
                latex_formulas['mode'] = r"Mode = \text{Most frequent value(s)}"

        return latex_formulas

    def __init__(self):
        self.formulas = {
            'mean': "Mean (Ungrouped) = Σ(x_i) / n",
            'median': "Median (Ungrouped) = Middle value or (x_{n/2} + x_{(n/2 + 1)}) / 2",
            'mode': "Mode (Ungrouped) = Most frequent value(s)"
        }

    def mct_word(self, filename="mct_data.docx"):
        """Converts statistical formulas into a Word document."""
        doc = Document()
        doc.add_heading('Measures of Central Tendency', 0)

        for name, formula in self.formulas.items():
            doc.add_heading(name.capitalize(), level=1)
            doc.add_paragraph(formula)

        doc.save(filename)
        print(f"Word document saved as {filename}")

    def mct_ppt(self, filename="mct_data.pptx"):
        """Creates a PowerPoint presentation with statistical formulas and animations."""
        ppt = Presentation()
        slide_layout = ppt.slide_layouts[1]

        slide = ppt.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.shapes.placeholders[1]

        title.text = "Measures of Central Tendency"
        content.text = "\n".join([f"{name.capitalize()}: {formula}" for name, formula in self.formulas.items()])

        ppt.save(filename)
        print(f"PowerPoint presentation saved as {filename}")

    def mct_exl(self, filename="mct_data.xlsx"):
        """Creates an Excel sheet with statistical formulas for tabular calculations."""
        wb = Workbook()
        ws = wb.active
        ws.title = "MCT Data"

        # Add headers
        ws.append(["Measure", "Formula"])

        # Add data
        for name, formula in self.formulas.items():
            ws.append([name.capitalize(), formula])

        wb.save(filename)
        print(f"Excel file saved as {filename}")

    def mct_pdf(self, filename="mct_data.pdf"):
        """Generates a PDF document with statistical formulas, including mathematical equation form."""
        pdf = FPDF()
        pdf.add_page()

        pdf.set_font("Arial", 'B', 16)
        pdf.cell(200, 10, "Measures of Central Tendency", ln=True, align="C")

        pdf.set_font("Arial", size=12)
        for name, formula in self.formulas.items():
            pdf.ln(10)
            pdf.cell(200, 10, f"{name.capitalize()}: {formula}", ln=True)

        pdf.output(filename)
        print(f"PDF saved as {filename}")

    def __init__(self):
        self.formulas = {
            'mean': "Mean (Ungrouped) = Σ(x_i) / n",
            'median': "Median (Ungrouped) = Middle value or (x_{n/2} + x_{(n/2 + 1)}) / 2",
            'mode': "Mode (Ungrouped) = Most frequent value(s)"
        }

    def mct_text(self, filename="mct_data.txt"):
        """Creates a plain text file with statistical formulas."""
        with open(filename, 'w') as file:
            file.write("Measures of Central Tendency\n")
            file.write("="*30 + "\n")
            for name, formula in self.formulas.items():
                file.write(f"{name.capitalize()}:\n")
                file.write(f"{formula}\n\n")
        print(f"Text file saved as {filename}")

    def mct_csv(self, filename="mct_data.csv"):
        """Creates a CSV file with statistical formulas."""
        with open(filename, 'w') as file:
            file.write("Measure,Formula\n")
            for name, formula in self.formulas.items():
                file.write(f"{name.capitalize()},{formula}\n")
        print(f"CSV file saved as {filename}")

    def mct_html(self, filename="mct_data.html"):
        """Creates an HTML file with statistical formulas."""
        with open(filename, 'w') as file:
            file.write("<html><head><title>MCT Data</title></head><body>\n")
            file.write("<h1>Measures of Central Tendency</h1>\n")
            for name, formula in self.formulas.items():
                file.write(f"<h2>{name.capitalize()}</h2>\n")
                file.write(f"<p>{formula}</p>\n")
            file.write("</body></html>\n")
        print(f"HTML file saved as {filename}")

    def mct_latx(self, filename="mct_data.tex"):
        """Creates a LaTeX file with statistical formulas."""
        with open(filename, 'w') as file:
            file.write(r"\documentclass{article}\begin{document}\n")
            file.write(r"\section*{Measures of Central Tendency}\n")
            for name, formula in self.formulas.items():
                file.write(rf"\subsection*{{{name.capitalize()}}}\n")
                file.write(rf"{formula} \\[2mm]\n")
            file.write(r"\end{document}")
        print(f"LaTeX file saved as {filename}")

    def __init__(self, data, grouped=False):
        """
        Initializes the MainRange class with the provided data.

        Parameters:
        - data: A list of numerical values or a dictionary for grouped data.
        - grouped: Boolean, if True, treats the data as grouped.
        """
        self.data = data
        self.grouped = grouped

    @staticmethod
    def range(data, grouped=False):
        """Calculates the range of the data."""
        if grouped:
            return max(data.keys()) - min(data.keys())
        else:
            return max(data) - min(data)

    @staticmethod
    def variance(data, grouped=False):
        """Calculates the variance of the data."""
        if grouped:
            total_frequency = sum(data.values())
            mean = sum(key * freq for key, freq in data.items()) / total_frequency
            return sum(freq * ((key - mean) ** 2) for key, freq in data.items()) / total_frequency
        else:
            mean = sum(data) / len(data)
            return sum((x - mean) ** 2 for x in data) / len(data)

    @staticmethod
    def std_dev(data, grouped=False):
        """Calculates the standard deviation of the data."""
        variance = MainRange.variance(data, grouped)
        return variance ** 0.5

    @staticmethod
    def iqr_range(data, grouped=False):
        """Calculates the interquartile range (IQR) of the data."""
        if grouped:
            sorted_data = []
            for key, freq in sorted(data.items()):
                sorted_data.extend([key] * freq)
        else:
            sorted_data = sorted(data)

        n = len(sorted_data)

        # Calculate Q1 (25th percentile)
        q1_index = (n + 1) // 4
        q1 = sorted_data[q1_index - 1]

        # Calculate Q3 (75th percentile)
        q3_index = (3 * (n + 1)) // 4
        q3 = sorted_data[q3_index - 1]

        return q3 - q1

    @staticmethod
    def mod(data, grouped=False):
        """
        Calculates the mode of the data.

        Returns:
        - A list of mode values. If no mode exists, returns an empty list.
        """
        if grouped:
            sorted_data = []
            for key, freq in sorted(data.items()):
                sorted_data.extend([key] * freq)
        else:
            sorted_data = data

        frequency = {}
        for value in sorted_data:
            if value in frequency:
                frequency[value] += 1
            else:
                frequency[value] = 1

        max_freq = max(frequency.values())
        modes = [value for value, freq in frequency.items() if freq == max_freq]

        return modes if max_freq > 1 else []

    @staticmethod
    def mod_wr(data, grouped=False):
        """
        Returns the working table for calculating the mode of the data.

        Returns:
        - A list of dictionaries representing the frequency table for ungrouped data 
          or a detailed frequency distribution for grouped data.
        """
        if grouped:
            frequency_table = [{"Value": key, "Frequency": freq} for key, freq in sorted(data.items())]
        else:
            frequency_table = {}
            for value in data:
                if value in frequency_table:
                    frequency_table[value] += 1
                else:
                    frequency_table[value] = 1
            
            # Convert frequency table to list of dictionaries
            frequency_table = [{"Value": key, "Frequency": freq} for key, freq in frequency_table.items()]

        return frequency_table

    @staticmethod
    def mod_formulae(grouped=False, return_all=True, return_variance=False, return_std_dev=False, return_iqr=False):
        """
        Returns the formulae for calculating mode, variance, standard deviation, and interquartile range (IQR)
        along with detailed descriptions.

        Parameters:
        - grouped: Boolean, if True, returns formulae for grouped data. Default is False (ungrouped data).
        - return_all: Boolean, if True, returns formulae for all measures. Default is True.
        - return_variance: Boolean, if True, returns the formula for variance only.
        - return_std_dev: Boolean, if True, returns the formula for standard deviation only.
        - return_iqr: Boolean, if True, returns the formula for interquartile range only.

        Returns:
        - A dictionary containing the formulae and their descriptions for the requested measures.
        """
        formulae = {}

        # Mode formula with description
        if return_all or (not grouped and not return_variance and not return_std_dev and not return_iqr):
            formulae['mode'] = {
                'formula': "Mode (Ungrouped) = value with the highest frequency",
                'description': "The mode is the value that appears most frequently in a dataset.",
                'parameters': {
                    'value': "The data point that appears most often."
                }
            }

        if return_all or (grouped and not return_variance and not return_std_dev and not return_iqr):
            formulae['mode_grouped'] = {
                'formula': "Mode (Grouped) = x_i where f_i is maximum",
                'description': "The mode for grouped data is the class interval with the highest frequency.",
                'parameters': {
                    'x_i': "Midpoint of the modal class.",
                    'f_i': "Frequency of the modal class."
                }
            }

        # Variance formula with description
        if return_all or return_variance:
            if grouped:
                formulae['variance_grouped'] = {
                    'formula': "Variance (Grouped) = Σ(f_i * (x_i - mean)^2) / Σf_i",
                    'description': "Variance for grouped data is calculated using the frequencies and midpoints.",
                    'parameters': {
                        'f_i': "Frequency of the i-th class.",
                        'x_i': "Midpoint of the i-th class.",
                        'mean': "Mean of the data.",
                        'Σf_i': "Total sum of frequencies."
                    }
                }
            else:
                formulae['variance'] = {
                    'formula': "Variance (Ungrouped) = Σ(x_i - mean)^2 / N",
                    'description': "Variance for ungrouped data is calculated using the individual data points.",
                    'parameters': {
                        'x_i': "Individual data point.",
                        'mean': "Mean of the data.",
                        'N': "Total number of observations."
                    }
                }

        # Standard deviation formula with description
        if return_all or return_std_dev:
            if grouped:
                formulae['std_dev_grouped'] = {
                    'formula': "Standard Deviation (Grouped) = √(Variance (Grouped))",
                    'description': "Standard deviation is the square root of variance for grouped data.",
                    'parameters': {}
                }
            else:
                formulae['std_dev'] = {
                    'formula': "Standard Deviation (Ungrouped) = √(Variance (Ungrouped))",
                    'description': "Standard deviation is the square root of variance for ungrouped data.",
                    'parameters': {}
                }

        # Interquartile Range formula with description
        if return_all or return_iqr:
            if grouped:
                formulae['iqr_grouped'] = {
                    'formula': "IQR (Grouped) = Q3 - Q1",
                    'description': "The interquartile range is the difference between the third (Q3) and first (Q1) quartiles.",
                    'parameters': {
                        'Q1': "First quartile (25th percentile).",
                        'Q3': "Third quartile (75th percentile)."
                    }
                }
            else:
                formulae['iqr'] = {
                    'formula': "IQR (Ungrouped) = Q3 - Q1",
                    'description': "The interquartile range is the difference between the third (Q3) and first (Q1) quartiles.",
                    'parameters': {
                        'Q1': "First quartile (25th percentile).",
                        'Q3': "Third quartile (75th percentile)."
                    }
                }

        return formulae
#. Examples 
if __name__ == "__main__":
    formulae = Hgdgdhdhsb.mod_formulae(grouped=True, return_all=True)
    for measure, details in formulae.items():
        print(f"{measure.capitalize()}:\nFormula: {details['formula']}\nDescription: {details['description']}\n")
        if details['parameters']:
            print("Parameters:")
            for param, desc in details['parameters'].items():
                print(f" - {param}: {desc}")
        print("\n")
#. End examples

    @staticmethod
    def histogram(data, bins=10, grouped=False, ax=None):
        """Creates a histogram representation of the data."""
        if grouped:
            frequencies = list(data.values())
            bin_keys = list(data.keys())
            min_value = min(bin_keys)
            max_value = max(bin_keys)
            bin_size = (max_value - min_value) / bins
            histogram_freq = [0] * bins

            for key in bin_keys:
                index = int((key - min_value) / bin_size)
                if index >= bins:
                    index = bins - 1
                histogram_freq[index] += data[key]

            ax.bar([min_value + i * bin_size for i in range(bins)],
                    histogram_freq,
                    width=bin_size,
                    edgecolor='black')
            ax.set_title('Histogram (Grouped)')
            ax.set_xlabel('Bins')
            ax.set_ylabel('Frequency')
        else:
            ax.hist(data, bins=bins, edgecolor='black')
            ax.set_title('Histogram (Ungrouped)')
            ax.set_xlabel('Value')
            ax.set_ylabel('Frequency')

    @staticmethod
    def box(data, grouped=False, ax=None):
        """Creates a box representation of the data."""
        if grouped:
            sorted_data = []
            for key in sorted(data.keys()):
                sorted_data.extend([key] * data[key])
        else:
            sorted_data = sorted(data)

        ax.boxplot(sorted_data)
        ax.set_title('Box')
        ax.set_ylabel('Values') 

    @staticmethod
    def scatter(x_data, y_data, grouped=False, ax=None):
        """Creates a scatter representation of the data."""
        if grouped:
            raise ValueError("Scatter plots are typically not suitable for grouped data.")
        if len(x_data) != len(y_data):
            raise ValueError("x_data and y_data must have the same length.")

        ax.scatter(x_data, y_data)
        ax.set_title('Scatter')
        ax.set_xlabel('X Data')
        ax.set_ylabel('Y Data')

    @staticmethod
    def pie(categories, values, grouped=False, ax=None):
        """Creates a pie representation of the data."""
        if grouped:
            ax.pie(categories.values(), labels=categories.keys(), autopct='%1.1f%%')
            ax.set_title('Pie (Grouped)')
        else:
            ax.pie(values, labels=categories, autopct='%1.1f%%')
            ax.set_title('Pie (Ungrouped)')

    @staticmethod
    def dv(data, x_data=None, y_data=None, categories=None, values=None, grouped=False, graph_options=None):
        """Displays selected graphs in one window based on user options."""
        fig, axs = plt.subplots(2, 2, figsize=(12, 10))
        axs = axs.flatten()  # Flatten the 2D array of axes for easy indexing
        
        # Default graph options if none are provided
        if graph_options is None:
            graph_options = ['histogram', 'box', 'scatter', 'pie']
        
        # Track which graphs to display
        graphs_displayed = 0

        if 'histogram' in graph_options:
            YYY.display_histogram(data, axs[graphs_displayed], grouped)
            graphs_displayed += 1

        if 'box' in graph_options:
            YYY.display_box(data, axs[graphs_displayed], grouped)
            graphs_displayed += 1

        if 'scatter' in graph_options:
            if x_data is not None and y_data is not None:
                YYY.display_scatter(x_data, y_data, axs[graphs_displayed], grouped)
                graphs_displayed += 1

        if 'pie' in graph_options:
            if categories is not None and values is not None:
                YYY.display_pie(categories, values, axs[graphs_displayed], grouped)
                graphs_displayed += 1

        # Remove unused axes
        for i in range(graphs_displayed, len(axs)):
            fig.delaxes(axs[i])

        plt.tight_layout()
        plt.show()

    @staticmethod
    def display_histogram(data, ax, grouped):
        """Helper function to display histogram."""
        staT.histogram(data, ax=ax, grouped=grouped)

    @staticmethod
    def display_box(data, ax, grouped):
        """Helper function to display box."""
        DV.box(data, ax=ax, grouped=grouped)

    @staticmethod
    def display_scatter(x_data, y_data, ax, grouped):
        """Helper function to display scatter."""
        if grouped:
            raise ValueError("Scatter representations are typically not suitable for grouped data.")
        staT.scatter(x_data, y_data, ax=ax, grouped=grouped)

    @staticmethod
    def display_pie(categories, values, ax, grouped):
        """Helper function to display pie."""
        staT.pie(categories, values, ax=ax, grouped=grouped)

# ___________________
# Usage example
if __name__ == "__main__":
    # Sample ungrouped data
    ungrouped_data = [1, 2, 2, 3, 3, 3, 4, 5, 5, 5, 6, 7, 8, 9]
    
    # Sample grouped data
    grouped_data = {
        1: 2,
        2: 3,
        3: 4,
        4: 1,
        5: 3,
        6: 1,
        7: 1,
        8: 1,
        9: 1,
    }
    
    categories = ['A', 'B', 'C', 'D']
    values = [10, 15, 7, 8]

    # Display selected graphs for ungrouped data
    print("Displaying Histogram and Box for Ungrouped Data:")
    YYY.dv(ungrouped_data, x_data=ungrouped_data, y_data=[x * 2 for x in ungrouped_data],
            categories=categories, values=values, grouped=False, graph_options=['histogram', 'box'])

    print("\n--- Grouped Data Visualizations ---\n")
    
    # Display all graphs for grouped data
    print("Displaying Pie and Histogram for Grouped Data:")
    YYY.dv(grouped_data, categories=grouped_data, grouped=True, graph_options=['pie', 'histogram'])  # Note: Scatter is not displayed for grouped data
 # _______________
 
    @staticmethod
    def details(methods=None):
        """Prints details about the available visualization methods."""
        all_methods = {
            "histogram": "Creates a histogram representation of the data.\n"
                         "Can be used for both grouped and ungrouped data.\n"
                         "Parameters: data (list or dict), bins (int), grouped (bool)",
            "box_plot": "Creates a box plot representation of the data.\n"
                        "Displays the minimum, first quartile (Q1), median, third quartile (Q3), and maximum values.\n"
                        "Can be used for both grouped and ungrouped data.\n"
                        "Parameters: data (list or dict), grouped (bool)",
            "scatter_plot": "Creates a scatter plot representation of the data.\n"
                            "Typically used for ungrouped data with paired x and y values.\n"
                            "Parameters: x_data (list), y_data (list), grouped (bool)",
            "pie_chart": "Creates a pie chart representation of the data.\n"
                         "Displays the proportion of each category.\n"
                         "Can be used for both grouped and ungrouped data.\n"
                         "Parameters: categories (list or dict), values (list), grouped (bool)"
        }

        if methods is None:
            methods = all_methods.keys()

        print("Selective Data Visualization Methods:")
        for method in methods:
            if method in all_methods:
                print(f"{method.capitalize()}:")
                print(all_methods[method])
                print()
            else:
                print(f"Method '{method}' not found.")
    
    @staticmethod
    def histogram_plot(data, bins=10, grouped=False):
        """Creates a histogram representation of the data."""
        if grouped:
            frequencies = list(data.values())
            bin_keys = list(data.keys())
            min_value = min(bin_keys)
            max_value = max(bin_keys)
            bin_size = (max_value - min_value) / bins
            histogram_freq = [0] * bins

            for key in bin_keys:
                index = int((key - min_value) / bin_size)
                if index >= bins:
                    index = bins - 1
                histogram_freq[index] += data[key]

            # Print histogram
            print("Histogram (Grouped):")
            for i in range(bins):
                lower_bound = min_value + i * bin_size
                upper_bound = min_value + (i + 1) * bin_size
                print(f"{lower_bound:.2f} - {upper_bound:.2f}: {'#' * histogram_freq[i]} ({histogram_freq[i]})")
        else:
            min_value = min(data)
            max_value = max(data)
            bin_size = (max_value - min_value) / bins
            frequencies = [0] * bins

            for value in data:
                index = int((value - min_value) / bin_size)
                if index >= bins:
                    index = bins - 1
                frequencies[index] += 1

            # Print histogram
            print("Histogram (Ungrouped):")
            for i in range(bins):
                print(f"{min_value + i * bin_size:.2f} - {min_value + (i + 1) * bin_size:.2f}: {'#' * frequencies[i]} ({frequencies[i]})")

    @staticmethod
    def box_plot(data, grouped=False):
        """Creates a box plot representation of the data."""
        if grouped:
            sorted_data = []
            for key in sorted(data.keys()):
                sorted_data.extend([key] * data[key])
        else:
            sorted_data = sorted(data)

        n = len(sorted_data)
        q1 = sorted_data[n // 4]
        median = sorted_data[n // 2]
        q3 = sorted_data[(3 * n) // 4]
        min_value = min(sorted_data)
        max_value = max(sorted_data)

        # Print box plot
        print("Box Plot:")
        print(f"Min: {min_value}")
        print(f"Q1: {q1}")
        print(f"Median: {median}")
        print(f"Q3: {q3}")
        print(f"Max: {max_value}")

    @staticmethod
    def scatter_plot(x_data, y_data, grouped=False):
        """Creates a scatter plot representation of the data."""
        if grouped:
            raise ValueError("Scatter plots are typically not suitable for grouped data.")
        if len(x_data) != len(y_data):
            raise ValueError("x_data and y_data must have the same length.")

        print("Scatter Plot:")
        for x, y in zip(x_data, y_data):
            print(f"({x}, {y})")

    @staticmethod
    def pie_chart(categories, values, grouped=False):
        """Creates a pie chart representation of the data."""
        total = sum(values)
        print("Pie Chart:")
        if grouped:
            for category, value in categories.items():
                percentage = (value / total) * 100
                print(f"{category}: {'#' * int(percentage // 2)} ({percentage:.1f}%)")
        else:
            for category, value in zip(categories, values):
                percentage = (value / total) * 100
                print(f"{category}: {'#' * int(percentage // 2)} ({percentage:.1f}%)")
